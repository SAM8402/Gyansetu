import fitz
from docx import Document as DocxDocument
from pptx import Presentation
import re
from pathlib import Path
from app.core.logging_config import logger

class DocumentIntelService:
    async def process(self, file_path: str, doc_type: str = "auto") -> dict:
        path = Path(file_path)
        ext = path.suffix.lower()
        logger.info("document_intel_start", file=path.name, ext=ext, doc_type=doc_type)
        if ext == ".pdf":
            return self._parse_pdf(path, doc_type=doc_type)
        elif ext == ".docx":
            return self._parse_docx(path, doc_type=doc_type)
        elif ext == ".pptx":
            return self._parse_pptx(path)
        else:
            return self._parse_text(path)

    def _parse_pdf(self, path: Path, doc_type: str = "auto") -> dict:
        doc = fitz.open(path)
        raw_text = ""
        sections = []
        tables = []
        figures = []
        equations = []
        skip_heavy = doc_type == "mostly_text"
        for page_num, page in enumerate(doc):
            text = page.get_text()
            raw_text += f"\n--- Page {page_num + 1} ---\n{text}"
            blocks = page.get_text("dict")["blocks"]
            for block in blocks:
                if block["type"] == 0:
                    for line in block["lines"]:
                        for span in line["spans"]:
                            font_size = span["size"]
                            text_span = span["text"].strip()
                            if not text_span:
                                continue
                            if font_size > 14:
                                sections.append({"heading": text_span, "content": "", "level": 1})
                            elif font_size > 12:
                                sections.append({"heading": text_span, "content": "", "level": 2})
                            else:
                                if sections:
                                    sections[-1]["content"] += text_span + " "
            if not skip_heavy and doc_type in ("auto", "tables"):
                tabs = page.find_tables()
                for tab in tabs:
                    tables.append({"caption": "", "data": tab.extract()})
            if not skip_heavy and doc_type in ("auto", "diagrams_equations", "scanned"):
                images = page.get_images()
                for img in images:
                    figures.append({"caption": "", "page": page_num + 1})
        if doc_type in ("auto", "diagrams_equations"):
            eq_pattern = r'\$[^$]+\$|\\\[.*?\\\]|\\\(.*?\\\)'
            equations = re.findall(eq_pattern, raw_text)
        page_count = len(doc)
        doc.close()
        return {"raw_text": raw_text, "sections": sections, "tables": tables, "figures": figures, "equations": equations, "metadata": {"page_count": page_count, "word_count": len(raw_text.split()), "file_type": "pdf", "parsed_strategy": doc_type}}

    def _parse_docx(self, path: Path) -> dict:
        doc = DocxDocument(path)
        raw_text = "\n".join(p.text for p in doc.paragraphs)
        sections = []
        for p in doc.paragraphs:
            if p.style.name.startswith("Heading"):
                level = int(p.style.name.split()[-1]) if p.style.name.split()[-1].isdigit() else 1
                sections.append({"heading": p.text, "content": "", "level": level})
            elif sections:
                sections[-1]["content"] += p.text + " "
        tables = []
        for t in doc.tables:
            data = [[cell.text for cell in row.cells] for row in t.rows]
            tables.append({"caption": "", "data": data})
        eq_pattern = r'\$[^$]+\$|\\\[.*?\\\]|\\\(.*?\\\)'
        equations = re.findall(eq_pattern, raw_text)
        return {"raw_text": raw_text, "sections": sections, "tables": tables, "figures": [], "equations": equations, "metadata": {"page_count": len(doc.paragraphs), "word_count": len(raw_text.split()), "file_type": "docx"}}

    def _parse_pptx(self, path: Path) -> dict:
        prs = Presentation(path)
        raw_text = ""
        sections = []
        for slide_num, slide in enumerate(prs.slides):
            raw_text += f"\n--- Slide {slide_num + 1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    raw_text += shape.text + "\n"
        return {"raw_text": raw_text, "sections": sections, "tables": [], "figures": [], "equations": [], "metadata": {"page_count": len(prs.slides), "word_count": len(raw_text.split()), "file_type": "pptx"}}

    def _parse_text(self, path: Path) -> dict:
        text = path.read_text(encoding="utf-8", errors="ignore")
        return {"raw_text": text, "sections": [], "tables": [], "figures": [], "equations": [], "metadata": {"page_count": 1, "word_count": len(text.split()), "file_type": "txt"}}
